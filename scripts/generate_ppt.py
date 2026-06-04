#!/usr/bin/env python3
"""PPT 슬라이드 생성 — 기존 양식 100% 재현"""
import sys, json, copy, os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

def make_rpr(sz=1100, bold=False, color=None, font="Malgun Gothic", hyperlink_id=None):
    rpr = etree.Element(qn("a:rPr"))
    rpr.set("lang", "ko-KR")
    rpr.set("altLang", "en-US")
    rpr.set("sz", str(sz))
    rpr.set("dirty", "0")
    if bold:
        rpr.set("b", "1")
    if color:
        sf = etree.SubElement(rpr, qn("a:solidFill"))
        clr = etree.SubElement(sf, qn("a:srgbClr"))
        clr.set("val", color)
    if hyperlink_id:
        hl = etree.SubElement(rpr, qn("a:hlinkClick"))
        hl.set("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id", hyperlink_id)
    lat = etree.SubElement(rpr, qn("a:latin"))
    lat.set("typeface", font)
    ea = etree.SubElement(rpr, qn("a:ea"))
    ea.set("typeface", font)
    return rpr

def make_run(text, sz=1100, bold=False, color=None, font="Malgun Gothic", hyperlink_id=None):
    r = etree.Element(qn("a:r"))
    r.append(make_rpr(sz, bold, color, font, hyperlink_id))
    t = etree.SubElement(r, qn("a:t"))
    t.text = text
    return r

def make_bullet_para(text, level=0):
    """레벨 0: 제목 불렛(•), 레벨 1: 본문(―)"""
    p = etree.Element(qn("a:p"))
    pPr = etree.SubElement(p, qn("a:pPr"))
    if level == 0:
        pPr.set("marL", "171450")
        pPr.set("indent", "-171450")
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", "맑은 고딕")
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "•")
    else:
        pPr.set("marL", "361950")
        pPr.set("lvl", "1")
        pPr.set("indent", "-180975")
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", "바탕")
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "―")
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
    spcPct.set("val", "100000" if level == 0 else "150000")
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
    spcPts.set("val", "600")
    etree.SubElement(pPr, qn("a:buClrTx"))
    p.append(make_run(text))
    return p

def make_impl_para(text):
    """☞ 파란색 제언"""
    p = etree.Element(qn("a:p"))
    pPr = etree.SubElement(p, qn("a:pPr"))
    pPr.set("marL", "180975")
    pPr.set("lvl", "1")
    pPr.set("indent", "0")
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
    spcPct.set("val", "150000")
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
    spcPts.set("val", "600")
    etree.SubElement(pPr, qn("a:buClrTx"))
    etree.SubElement(pPr, qn("a:buNone"))
    p.append(make_run(f"☞  {text}", sz=1100, color="0070C0"))
    return p

def make_source_para(src_name, src_title, src_url, date, rel_id):
    """출처 줄 — 하이퍼링크 포함"""
    p = etree.Element(qn("a:p"))
    pPr = etree.SubElement(p, qn("a:pPr"))
    pPr.set("marL", "180975")
    pPr.set("lvl", "1")
    pPr.set("indent", "0")
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
    spcPct.set("val", "150000")
    etree.SubElement(pPr, qn("a:buClrTx"))
    etree.SubElement(pPr, qn("a:buNone"))
    p.append(make_run(f"출처: {src_name}, ", sz=900, color="000000", font="+mn-ea"))
    p.append(make_run(src_title, sz=900, color="0563C1", font="+mn-ea", hyperlink_id=rel_id))
    p.append(make_run(f" ({date})", sz=900, color="000000", font="+mn-ea"))
    return p

def generate_slide(data: dict, template_path: str, output_path: str):
    prs = Presentation(template_path)
    slide = prs.slides[0]

    country_map = {"JP": "일본", "US": "미국", "EU": "유럽", "AU": "호주"}

    # 국가 레이블 텍스트프레임 찾기 (Title 1)
    for shape in slide.shapes:
        if shape.name == "Title 1":
            tf = shape.text_frame
            tf.paragraphs[0].runs[0].text = country_map.get(data["country"], data["country"])

    # 본문 텍스트프레임 찾기 (idx=15)
    for shape in slide.shapes:
        if hasattr(shape, "placeholder_format") and shape.placeholder_format and shape.placeholder_format.idx == 15:
            tf = shape.text_frame
            txBody = tf._txBody
            # 기존 단락 전부 제거
            for p in txBody.findall(qn("a:p")):
                txBody.remove(p)

            # 제목 불렛
            txBody.append(make_bullet_para(f"{data['title']} ({data['date']})", level=0))

            # 본문 불렛
            for bullet in data.get("bullets", []):
                txBody.append(make_bullet_para(bullet, level=1))

            # 제언
            txBody.append(make_impl_para(data.get("implication", "")))

            # 출처 하이퍼링크 rel 추가
            slide_part = slide.part
            rel = slide_part.relate_to(
                data["url"],
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
                is_external=True,
            )
            txBody.append(make_source_para(data["srcName"], data["title"], data["url"], data["date"], rel))
            break

    prs.save(output_path)
    print(f"Saved: {output_path}")

if __name__ == "__main__":
    payload = json.loads(sys.argv[1])
    template = sys.argv[2]
    output = sys.argv[3]
    generate_slide(payload, template, output)
